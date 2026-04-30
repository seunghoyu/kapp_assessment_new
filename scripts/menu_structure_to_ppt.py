#!/usr/bin/env python3
"""
menu-structure.planning.json → PPT 메뉴 구조도

- 기본 동작: 1뎁스(루트)마다 슬라이드 1장
- 옵션(--single-slide): 전체 루트를 한 슬라이드에 모두 배치

의존성:
  pip install python-pptx

실행 예:
  python scripts/menu_structure_to_ppt.py \\
    --input menu-structure.planning.json \\
    --output menu-structure.planning.pptx

검증 포인트:
  - PowerPoint에서 노드 겹침이 심하면 --margin 을 줄이거나 슬라이드를 수동 확대
  - 연결선은 부모 박스 하단 중앙 → 자식 박스 상단 중앙
"""

from __future__ import annotations

import argparse
import json
import math
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Inches, Pt
except ImportError as e:
    print("python-pptx 가 필요합니다: pip install python-pptx", file=sys.stderr)
    raise e


# depth(1~) → 배경색 (연한 톤, 텍스트는 진한 회색)
DEPTH_BG: list[tuple[int, int, int]] = [
    (219, 234, 254),
    (220, 252, 231),
    (254, 249, 195),
    (243, 232, 255),
    (255, 228, 230),
    (207, 250, 254),
    (229, 231, 235),
    (254, 243, 199),
    (233, 213, 255),
    (209, 250, 229),
    (253, 230, 138),
    (224, 242, 254),
]
TEXT_COLOR = RGBColor(17, 24, 39)
EDGE_COLOR = RGBColor(100, 116, 139)

EMU_PER_INCH = 914400


def _to_emu(v: float) -> int:
    return int(round(v * EMU_PER_INCH))


@dataclass
class TreeNode:
    메뉴명: str
    경로: str
    depth: int
    children: list[TreeNode] = field(default_factory=list)
    # layout
    subtree_width: float = 1.0
    subtree_left: float = 0.0
    center_x: float = 0.0


def parse_nodes(raw: dict[str, Any]) -> TreeNode:
    kids = raw.get("children") or []
    return TreeNode(
        메뉴명=str(raw.get("메뉴명", "")),
        경로=str(raw.get("경로", "")),
        depth=int(raw.get("depth", 1)),
        children=[parse_nodes(c) for c in kids],
    )


def layout_subtree(n: TreeNode, left: float, sibling_gap: float) -> float:
    """Returns subtree width in abstract X units (leaf = 1)."""
    if not n.children:
        n.subtree_left = left
        n.subtree_width = 1.0
        n.center_x = left + 0.5
        return 1.0

    cur = left
    for i, c in enumerate(n.children):
        if i > 0:
            cur += sibling_gap
        w = layout_subtree(c, cur, sibling_gap)
        cur += w

    total_w = cur - left
    n.subtree_left = left
    n.subtree_width = total_w
    n.center_x = left + total_w / 2.0
    return total_w


def layout_forest(roots: list[TreeNode], root_gap: float, sibling_gap: float) -> None:
    cur = 0.0
    for i, r in enumerate(roots):
        if i > 0:
            cur += root_gap
        layout_subtree(r, cur, sibling_gap)
        cur += r.subtree_width


def collect_nodes(n: TreeNode) -> list[TreeNode]:
    out = [n]
    for c in n.children:
        out.extend(collect_nodes(c))
    return out


def depth_bg(depth: int) -> tuple[int, int, int]:
    idx = max(0, depth - 1) % len(DEPTH_BG)
    return DEPTH_BG[idx]


def _render_slide_for_roots(
    prs: Presentation,
    roots: list[TreeNode],
    *,
    slide_w_in: float,
    slide_h_in: float,
    margin_in: float,
    sibling_gap: float,
    root_gap: float,
    row_step: float,
    min_font_pt: float,
    max_font_pt: float,
    title_text: str | None = None,
) -> None:
    layout_forest(roots, root_gap=root_gap, sibling_gap=sibling_gap)
    nodes: list[TreeNode] = []
    for r in roots:
        nodes.extend(collect_nodes(r))

    if not nodes:
        return

    depths = [n.depth for n in nodes]
    min_d, max_d = min(depths), max(depths)
    row_count = max(1, max_d - min_d + 1)

    cx = [n.center_x for n in nodes]
    half = 0.5
    min_x, max_x = min(cx) - half, max(cx) + half
    span_x = max(max_x - min_x, 1e-6)

    def row_index(d: int) -> int:
        return d - min_d

    span_y = max(row_count * row_step, 1e-6)

    title_h_in = 0.45 if title_text else 0.0
    usable_w = slide_w_in - 2 * margin_in
    usable_h = slide_h_in - 2 * margin_in - title_h_in
    scale = min(usable_w / span_x, usable_h / span_y)

    box_w = min(2.2 * scale, usable_w / max(12, math.sqrt(len(nodes))))
    box_h = box_w * 0.42
    box_w = max(box_w, 0.35)
    box_h = max(box_h, 0.18)

    font_pt = max(min_font_pt, min(max_font_pt, 7.5 * math.sqrt(scale)))

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    if title_text:
        title = slide.shapes.add_textbox(
            Inches(margin_in),
            Inches(margin_in * 0.6),
            Inches(slide_w_in - 2 * margin_in),
            Inches(title_h_in),
        )
        tf_t = title.text_frame
        tf_t.clear()
        p = tf_t.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR

    geom: dict[int, tuple[float, float, float, float]] = {}

    def iter_parent_child():
        for root in roots:
            q = [root]
            while q:
                p = q.pop(0)
                for c in p.children:
                    yield p, c
                    q.append(c)

    y_offset_in = margin_in + title_h_in

    for n in nodes:
        x_in = margin_in + (n.center_x - min_x) * scale - box_w / 2
        y_in = y_offset_in + row_index(n.depth) * row_step * scale
        geom[id(n)] = (x_in, y_in, box_w, box_h)

    for p, c in iter_parent_child():
        px, py, pw, ph = geom[id(p)]
        cx2, cy2, cw, ch = geom[id(c)]
        x1 = px + pw / 2
        y1 = py + ph
        x2 = cx2 + cw / 2
        y2 = cy2
        cx_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _to_emu(x1),
            _to_emu(y1),
            _to_emu(x2),
            _to_emu(y2),
        )
        cx_line.line.color.rgb = EDGE_COLOR
        cx_line.line.width = Pt(0.75)

    for n in nodes:
        x_in, y_in, bw, bh = geom[id(n)]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_in),
            Inches(y_in),
            Inches(bw),
            Inches(bh),
        )
        shape.fill.solid()
        r, g, b = depth_bg(n.depth)
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.color.rgb = EDGE_COLOR
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Emu(72000)
        tf.margin_right = Emu(72000)
        tf.margin_top = Emu(36000)
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = n.메뉴명
        p0.font.size = Pt(font_pt)
        p0.font.bold = True
        p0.font.color.rgb = TEXT_COLOR

        if n.경로.strip():
            p1 = tf.add_paragraph()
            p1.text = n.경로
            p1.level = 0
            p1.font.size = Pt(max(min_font_pt, font_pt - 2))
            p1.font.color.rgb = RGBColor(55, 65, 81)


def build_ppt(
    roots: list[TreeNode],
    out_path: Path,
    slide_w_in: float,
    slide_h_in: float,
    margin_in: float,
    sibling_gap: float,
    root_gap: float,
    row_step: float,
    min_font_pt: float,
    max_font_pt: float,
    *,
    per_root_slide: bool,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)

    if per_root_slide:
        for r in roots:
            _render_slide_for_roots(
                prs,
                [r],
                slide_w_in=slide_w_in,
                slide_h_in=slide_h_in,
                margin_in=margin_in,
                sibling_gap=sibling_gap,
                root_gap=root_gap,
                row_step=row_step,
                min_font_pt=min_font_pt,
                max_font_pt=max_font_pt,
                title_text=r.메뉴명,
            )
    else:
        _render_slide_for_roots(
            prs,
            roots,
            slide_w_in=slide_w_in,
            slide_h_in=slide_h_in,
            margin_in=margin_in,
            sibling_gap=sibling_gap,
            root_gap=root_gap,
            row_step=row_step,
            min_font_pt=min_font_pt,
            max_font_pt=max_font_pt,
            title_text="메뉴 구조도",
        )

    prs.save(str(out_path))


def main() -> None:
    ap = argparse.ArgumentParser(description="메뉴 JSON → PPT 트리")
    ap.add_argument(
        "--input",
        "-i",
        type=Path,
        default=Path("menu-structure.planning.json"),
        help="입력 JSON 경로",
    )
    ap.add_argument(
        "--output",
        "-o",
        type=Path,
        default=Path("menu-structure.planning.pptx"),
        help="출력 PPTX 경로",
    )
    ap.add_argument("--slide-w", type=float, default=13.333, help="슬라이드 너비(in), 기본 16:9")
    ap.add_argument("--slide-h", type=float, default=7.5, help="슬라이드 높이(in)")
    ap.add_argument("--margin", type=float, default=0.35, help="여백(in)")
    ap.add_argument("--sibling-gap", type=float, default=0.35, help="형제 노드 간 X 간격(추상 단위)")
    ap.add_argument("--root-gap", type=float, default=0.9, help="최상위 트리 간 X 간격(추상 단위)")
    ap.add_argument("--row-step", type=float, default=1.25, help="깊이 한 단계당 Y 간격(추상 단위)")
    ap.add_argument("--min-font", type=float, default=5.0, help="최소 글자 크기(pt)")
    ap.add_argument("--max-font", type=float, default=11.0, help="최대 글자 크기(pt)")
    ap.add_argument(
        "--single-slide",
        action="store_true",
        help="전체 루트를 한 슬라이드에 배치 (기본은 1뎁스마다 슬라이드 1장)",
    )
    args = ap.parse_args()

    raw = json.loads(args.input.read_text(encoding="utf-8"))
    if not isinstance(raw, list):
        raise SystemExit("JSON 루트는 배열이어야 합니다.")

    roots = [parse_nodes(item) for item in raw]

    build_ppt(
        roots,
        args.output,
        slide_w_in=args.slide_w,
        slide_h_in=args.slide_h,
        margin_in=args.margin,
        sibling_gap=args.sibling_gap,
        root_gap=args.root_gap,
        row_step=args.row_step,
        min_font_pt=args.min_font,
        max_font_pt=args.max_font,
        per_root_slide=(not args.single_slide),
    )
    print(f"Wrote {args.output.resolve()}")


if __name__ == "__main__":
    main()
